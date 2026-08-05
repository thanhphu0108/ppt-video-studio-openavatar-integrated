from __future__ import annotations

import io
import re
from collections import Counter
from typing import Any

from pptx import Presentation

DECORATIVE = {"MỞ ĐẦU", "QUY TRÌNH", "TÀI LIỆU", "CẢI TIẾN", "THAM KHẢO"}


def _norm(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _is_url(text: str) -> bool:
    return bool(re.match(r"^(https?://|www\.)", text, flags=re.I))


def _is_decorative(text: str) -> bool:
    value = _norm(text)
    return not value or value.upper() in DECORATIVE or bool(re.fullmatch(r"[\d\s|•\-–—.=]+", value))


def read_pptx(payload: bytes) -> list[dict[str, Any]]:
    prs = Presentation(io.BytesIO(payload))
    records: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        values: list[str] = []
        title = ""
        if slide.shapes.title and getattr(slide.shapes.title, "text", None):
            title = _norm(slide.shapes.title.text)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = _norm(getattr(shape, "text", ""))
            if text and text not in values:
                values.append(text)
        records.append({"slide": index, "title": title, "all_text": values})

    counts = Counter()
    originals: dict[str, str] = {}
    for rec in records:
        seen = set()
        for value in rec["all_text"]:
            key = value.casefold()
            originals[key] = value
            seen.add(key)
        counts.update(seen)
    threshold = max(2, int(len(records) * 0.30 + 0.999))
    repeated = {key for key, count in counts.items() if count >= threshold}

    cleaned: list[dict[str, Any]] = []
    for rec in records:
        title = _norm(rec["title"])
        bullets: list[str] = []
        links: list[str] = []
        for value in rec["all_text"]:
            value = _norm(value)
            if value.casefold() in repeated and value.casefold() != title.casefold():
                continue
            if value.casefold() == title.casefold():
                continue
            if _is_url(value):
                links.append(value)
                continue
            if _is_decorative(value):
                continue
            if value not in bullets:
                bullets.append(value)
        if not title or _is_decorative(title):
            title = bullets.pop(0) if bullets else f"Slide {rec['slide']}"
        word_count = len(re.findall(r"\S+", " ".join(bullets)))
        cleaned.append({
            "slide": rec["slide"],
            "title": title[:180],
            "bullets": bullets[:10],
            "links": links,
            "word_count": word_count,
        })
    total = len(cleaned)
    for rec in cleaned:
        rec["slide_type"] = classify_slide(rec, total)
    return cleaned


def classify_slide(record: dict[str, Any], total: int) -> str:
    idx = int(record["slide"])
    title = record["title"]
    bullets = record["bullets"]
    text = " ".join([title, *bullets]).upper()
    if idx == 1 and len(bullets) <= 3:
        return "Trang bìa"
    if record["links"] and record["word_count"] <= 12:
        return "Slide liên kết/video"
    if any(k in text for k in ["TÀI LIỆU THAM KHẢO", "REFERENCES", "THAM KHẢO"]):
        return "Tài liệu tham khảo"
    if any(k in text for k in ["KẾT LUẬN", "TÓM LẠI", "TỔNG KẾT", "TRÂN TRỌNG CẢM ƠN"]):
        return "Kết luận"
    if len(bullets) <= 2 and record["word_count"] <= 20 and idx not in (1, total):
        return "Phân phần"
    if any(k in text for k in ["QUY TRÌNH", "BƯỚC", "GIAI ĐOẠN", "WORKFLOW"]):
        return "Quy trình"
    if re.search(r"\d+([,.]\d+)?\s*(%|CA|HỒ SƠ|NGÀY|LƯỢT)", text):
        return "Số liệu"
    return "Nội dung"


def build_narration(record: dict[str, Any], style: str = "Thuyết trình chuyên nghiệp") -> str:
    title = record["title"].rstrip(". ")
    bullets = [b.rstrip(". ") for b in record["bullets"][:5]]
    kind = record["slide_type"]
    if kind == "Trang bìa":
        return f"Xin chào quý anh chị. Nội dung trình bày hôm nay là {title}. " + ("Chúng ta sẽ lần lượt đi qua các nội dung chính của chuyên đề." if not bullets else f"Trọng tâm gồm {', '.join(bullets[:3])}.")
    if kind == "Phân phần":
        return f"Tiếp theo, chúng ta chuyển sang phần {title}. " + (f"Phần này tập trung vào {', '.join(bullets[:2])}." if bullets else "")
    if kind == "Kết luận":
        return f"Tóm lại, {title}. " + " ".join(f"{b}." for b in bullets)
    if kind == "Tài liệu tham khảo":
        return f"Phần {title} cung cấp các tài liệu và nguồn tham khảo để tiếp tục tra cứu sau buổi trình bày."
    if kind == "Slide liên kết/video":
        return f"Slide này giới thiệu một liên kết hoặc tài nguyên minh họa liên quan đến {title}."
    if style == "Đọc gần nguyên bản":
        return f"{title}. " + " ".join(f"{b}." for b in bullets)
    if style == "Điều hành":
        return f"Về {title}, cần lưu ý: " + ("; ".join(bullets[:4]) + "." if bullets else "đây là nội dung cần được theo dõi và triển khai phù hợp.")
    return f"Ở nội dung {title}, " + ("các điểm chính gồm " + "; ".join(bullets[:4]) + "." if bullets else "chúng ta cần tập trung vào thông điệp trọng tâm của slide.")
